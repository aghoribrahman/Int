from django.shortcuts import render
import openai
from pptx import Presentation
import pandas as pd
from django.shortcuts import render
from ppt.forms import ExcelUploadForm
from ppt.models import ExcelFile
import json
from django.conf import settings
import os
from gtts import gTTS
from pptx.util import Inches, Pt


def generate_short_paragraph(prompt):
    # Set up your OpenAI API credentials
    openai.api_key = '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'  # Replace with your OpenAI API key

    # Make an API call to generate the short paragraph
    response = openai.Completion.create(
        engine='text-davinci-003',  # Choose the appropriate OpenAI engine
        prompt=prompt,
        max_tokens=400,  # Adjust the desired length of the generated paragraph
        n=1,
        stop=None,
        temperature=0.7,  # Adjust the temperature for controlling randomness
    )

    short_paragraph = response.choices[0].text.strip()
    return short_paragraph

def text_to_audio(paragraphs):
    for i, paragraph in enumerate(paragraphs):
        tts = gTTS(text=paragraph, lang='en')
        return tts
        #filename = f"output_{i}.mp3"
        #file_path = os.path.join('media', filename)
        #retutts.save(file_path)

def upload_excel(request):
    if request.method == 'POST':
        form = ExcelUploadForm(request.POST, request.FILES)
        if form.is_valid():
            excel_file = request.FILES['excel_file']
            
            df = pd.read_excel(excel_file)
            prompts = df['Prompt'].tolist()  # Replace 'Prompt Column Name' with the actual column name in your Excel sheet
            prompts = prompts
            paragraphs = []           
            for promp in prompts:
                short_paragraph = generate_short_paragraph(promp)
                paragraphs.append(short_paragraph)
          
            text_to_audio(paragraphs)
            create_ppt_slides(paragraphs)
            return render(request, 'ppt.html',{'value':short_paragraph})
    else:
        form = ExcelUploadForm()
    return render(request, 'ppt.html', {'form': form})

def create_ppt_slides(paragraphs):
    prs = Presentation()
    for i in paragraphs:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = 'Short Paragraph'
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(14) 
        text_frame = slide.placeholders[1].text_frame
        p = text_frame.add_paragraph()
        p.text = i
        p.font.size = Pt(12)
     
        audio_file = text_to_audio(i)
        slide.shapes.add_movie(audio_file, left=0, top=0, width=1, height=1)
    prs.save('output.pptx')  # Save the PowerPoint file

def generate_ppt_from_excel(prompts):
    prompt = prompts
    paragraphs = []
    for promp in prompt:
        short_paragraph = generate_short_paragraph(promp)
        paragraphs.append(short_paragraph)
    create_ppt_slides(paragraphs)
'''

'''